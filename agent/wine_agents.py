"""
Merged agent logic, utility functions, models, and services from the original app_2 codebase.
"""

# =============================
# MODELS
# =============================
import os
from dataclasses import dataclass
from dotenv import load_dotenv
from enum import Enum
from pydantic import BaseModel
from typing import List, Optional, Dict, Any

load_dotenv()

@dataclass
class WineRAGConfig:
    openai_api_key: str = os.getenv("OPENAI_API_KEY")
    pinecone_api_key: str = os.getenv("PINECONE_API_KEY")
    pinecone_environment: str = os.getenv("PINECONE_ENVIRONMENT")
    index_name: str = "wine-knowledge"
    embedding_model: str = "text-embedding-3-large"
    chunk_size: int = 3000
    chunk_overlap: int = 100

class WineAttribute(Enum):
    TASTE = "taste"
    STYLE = "style"
    BODY = "body"
    SWEETNESS = "sweetness"
    FOOD_PAIRING = "food_pairing"
    ALTERNATIVES = "alternatives"
    COLOR = "color"
    REGION = "region"

class QueryRequest(BaseModel):
    question: str

class QueryResponse(BaseModel):
    answer: str

# =============================
# UTILS
# =============================
# --- chunking_utils.py ---
from langchain.text_splitter import RecursiveCharacterTextSplitter

def chunk_text(text: str, chunk_size: int = 3000, chunk_overlap: int = 100):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)

# --- ocr_utils.py ---
import pytesseract
from PIL import Image, ImageOps
import numpy as np
import io
import logging

def preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    gray = ImageOps.grayscale(image)
    if min(gray.size) < 200:
        scale = 200 / min(gray.size)
        new_size = (int(gray.size[0]*scale), int(gray.size[1]*scale))
        gray = gray.resize(new_size, Image.LANCZOS)
    bw = gray.point(lambda x: 0 if x < 128 else 255, '1')
    return bw

def extract_text_from_image(image_bytes: bytes) -> str:
    try:
        image = Image.open(io.BytesIO(image_bytes))
        processed = preprocess_image_for_ocr(image)
        text = pytesseract.image_to_string(processed)
        return text.strip()
    except Exception as e:
        logging.warning(f"OCR extraction failed: {e}")
        return ""

# --- vision_utils.py ---
import openai
import base64

def extract_vision_content_from_image(image_bytes: bytes, openai_api_key: str) -> str:
    """
    Use OpenAI GPT-4o Vision API to extract tables, chart summaries, and visual info from an image.
    """
    base64_img = base64.b64encode(image_bytes).decode()
    client = openai.OpenAI(api_key=openai_api_key)
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a document analysis assistant. Extract all possible structured information, tables, lists, and key facts that could be relevant to answering questions. Be exhaustive and explicit."},
                {"role": "user", "content": [
                    {"type": "text", "text": "Extract and summarize all tables, charts, diagrams, and visual information from this image. If it's a table, output as Markdown. If it's a chart, summarize all key points and minor details. If it's a diagram or list, extract all relationships and entities. Otherwise, provide a detailed, descriptive caption. Be explicit and comprehensive so that a user could answer questions about the image using only your output."},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_img}"}}
                ]}
            ],
            max_tokens=1000
        )
        vision_content = response.choices[0].message.content
        return vision_content.strip()
    except Exception as e:
        logging.warning(f"Vision API failed: {e}")
        return ""

# --- vector_utils.py ---
import pinecone
from langchain_openai import OpenAIEmbeddings
from pinecone import ServerlessSpec
import threading

_vector_store_instance = None
_vector_store_lock = threading.Lock()

def get_vector_store(config):
    """
    Returns a singleton instance of VectorStore per process.
    Thread-safe for FastAPI usage.
    """
    global _vector_store_instance
    if _vector_store_instance is None:
        with _vector_store_lock:
            if _vector_store_instance is None:
                _vector_store_instance = VectorStore(config)
    return _vector_store_instance

class VectorStore:
    def __init__(self, config):
        self.config = config
        self.embeddings = OpenAIEmbeddings(
            openai_api_key=config.openai_api_key,
            model=config.embedding_model
        )
        self.setup_pinecone()

    def setup_pinecone(self):
        try:
            self.pc = pinecone.Pinecone(api_key=self.config.pinecone_api_key)
            if self.config.index_name not in self.pc.list_indexes().names():
                self.pc.create_index(
                    name=self.config.index_name,
                    dimension=3072,
                    metric="cosine",
                    spec=ServerlessSpec(
                        cloud="aws",
                        region="us-east-1"
                    )
                )
            self.index = self.pc.Index(self.config.index_name)
        except Exception as e:
            logging.error(f"Pinecone setup failed: {e}")
            raise

    def upsert_documents(self, documents):
        vectors = []
        for i, doc in enumerate(documents):
            embedding = self.embeddings.embed_query(doc['metadata']['text'])
            metadata = {k: str(v) for k, v in doc['metadata'].items()}
            vectors.append({
                'id': f"wine_doc_{i}",
                'values': embedding,
                'metadata': metadata
            })
        batch_size = 100
        for i in range(0, len(vectors), batch_size):
            batch = vectors[i:i + batch_size]
            self.index.upsert(vectors=batch)

    def query(self, query_text, top_k=5):
        import time
        try:
            t0 = time.perf_counter()
            query_embedding = self.embeddings.embed_query(query_text)
            t1 = time.perf_counter()
            results = self.index.query(
                vector=query_embedding,
                top_k=top_k,
                include_metadata=True
            )
            t2 = time.perf_counter()
            logging.info(f"[PROFILE] Embedding time: {t1 - t0:.3f}s, Pinecone query time: {t2 - t1:.3f}s")
            return results.get('matches', [])
        except Exception as e:
            logging.error(f"Vector search failed: {e}")
            return []

# =============================
# SERVICES
# =============================
from fastapi import UploadFile

# Forward declare agent_graph for reference
# (agent_graph and agent nodes will be merged below)

def process_document(file_path, openai_api_key, chunk_size=2500, chunk_overlap=100):
    ext = os.path.splitext(file_path)[1].lower()
    text_splitter = lambda text: chunk_text(text, chunk_size, chunk_overlap)
    if ext == '.pdf':
        return process_pdf(file_path, text_splitter)
    elif ext == '.pptx':
        return process_pptx(file_path, openai_api_key, text_splitter)
    else:
        raise ValueError("Unsupported file type. Only PDF and PPTX are supported.")

from pptx import Presentation
from langchain.schema import Document

def process_pdf(file_path, text_splitter):
    from langchain_community.document_loaders import PyPDFLoader
    loader = PyPDFLoader(file_path)
    documents = loader.load()
    full_text = "\n".join([doc.page_content for doc in documents])
    chunks = text_splitter(full_text)
    wine_documents = [
        {
            'metadata': {
                'text': chunk,
                'chunk_id': f"chunk_{i}",
                'source': os.path.basename(file_path),
                'doc_type': 'pdf',
            }
        } for i, chunk in enumerate(chunks)
    ]
    return wine_documents

def process_pptx(file_path, openai_api_key, text_splitter):
    prs = Presentation(file_path)
    slide_documents = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        seen_ocr = set()
        for shape_idx, shape in enumerate(slide.shapes):
            # Extract text from shapes
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_text.append(shape.text.strip())
            # Extract from images (OCR + Vision)
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                ocr_text = extract_text_from_image(image_bytes)
                if ocr_text:
                    slide_text.append(ocr_text)
                vision_text = extract_vision_content_from_image(image_bytes, openai_api_key)
                if vision_text:
                    slide_text.append(vision_text)
        # Deduplicate and filter
        slide_text = [t for t in set(slide_text) if len(t.strip()) > 3]
        slide_text_combined = "\n".join(slide_text)
        chunks = text_splitter(slide_text_combined)
        for i, chunk in enumerate(chunks):
            slide_documents.append({
                'metadata': {
                    'text': chunk,
                    'chunk_id': f"slide_{slide_idx+1}_chunk_{i}",
                    'slide': slide_idx+1,
                    'source': os.path.basename(file_path),
                    'doc_type': 'pptx',
                }
            })
    return slide_documents

class AgenticRAGService:
    def __init__(self):
        self.config = WineRAGConfig()
        if not all([
            self.config.openai_api_key,
            self.config.pinecone_api_key,
            self.config.pinecone_environment
        ]):
            raise RuntimeError("Missing required environment variables.")
        self.vector_store = get_vector_store(self.config)
        self.agent_graph = build_agent_graph()

    def query(self, question: str):
        initial_state = {"query": question}
        result = self.agent_graph.invoke(initial_state)
        return QueryResponse(answer=result.get("final_response", "Sorry, I couldn't process your query."))

    def process_document(self, file: UploadFile):
        import tempfile
        # Save uploaded file to temp
        suffix = os.path.splitext(file.filename)[1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file.file.read())
            tmp_path = tmp.name
        try:
            docs = process_document(tmp_path, self.config.openai_api_key, self.config.chunk_size, self.config.chunk_overlap)
            self.vector_store.upsert_documents(docs)
            status = "processed"
        except Exception as e:
            status = f"error: {e}"
        finally:
            os.remove(tmp_path)
        return {"filename": file.filename, "status": status}

# =============================
# AGENT LOGIC (NODES & GRAPH)
# =============================
from langgraph.graph import StateGraph, END

def analyze_query_node(state: Dict) -> Dict[str, Any]:
    import openai
    config = WineRAGConfig()
    query = state["query"]
    prompt = f"""
    Analyze this wine-related query and extract information:
    Query: "{query}"
    Identify and return JSON with:
    1. "attributes": wine characteristics mentioned (body, sweetness, color, taste, etc.)
    2. "wines": specific wine names mentioned
    3. "query_type": one of [recommendation, comparison, information, pairing, alternatives]
    4. "foods": food items mentioned
    5. "preferences": any explicit preferences (e.g. dry, sweet, full-bodied, etc.)
    Output JSON only.
    """
    try:
        client = openai.OpenAI(api_key=config.openai_api_key)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
            max_tokens=500
        )
        import json
        analysis = json.loads(response.choices[0].message.content)
    except Exception as e:
        analysis = {"attributes": [], "wines": [], "query_type": "information", "foods": [], "preferences": []}
    return {"query_analysis": analysis, "query": query}

def vector_search_node(state: Dict) -> Dict[str, Any]:
    config = WineRAGConfig()
    vs = get_vector_store(config)
    query = state["query"]
    matches = vs.query(query, top_k=5)
    retrieved_chunks = [m['metadata'].get('text', '') for m in matches if m.get('metadata')]
    text_sufficient = bool(retrieved_chunks)
    next_node = "answer_synthesis" if text_sufficient else "vision_analysis"
    return {
        "query": state["query"],
        "retrieved_chunks": retrieved_chunks,
        "text_sufficient": text_sufficient,
        "vector_matches": matches,
        "__next__": next_node
    }

def vision_analysis_node(state: Dict) -> Dict[str, Any]:
    import time
    t0 = time.perf_counter()
    vision_chunks = []
    for match in state.get("vector_matches", []):
        if match.get('metadata', {}).get('doc_type') == 'pptx':
            vision_chunks.append(match['metadata'].get('text', ''))
    vision_sufficient = bool(vision_chunks)
    t1 = time.perf_counter()
    logging.info(f"[PROFILE] Vision analysis time: {t1 - t0:.3f}s")
    next_node = "answer_synthesis" if vision_sufficient else "ocr_analysis"
    return {
        "query": state["query"],
        "vision_chunks": vision_chunks,
        "vision_sufficient": vision_sufficient,
        "__next__": next_node
    }

def ocr_analysis_node(state: Dict) -> Dict[str, Any]:
    import time
    t0 = time.perf_counter()
    ocr_chunks = []
    for match in state.get("vector_matches", []):
        if match.get('metadata', {}).get('doc_type') == 'pptx':
            ocr_chunks.append(match['metadata'].get('text', ''))
    t1 = time.perf_counter()
    logging.info(f"[PROFILE] OCR analysis time: {t1 - t0:.3f}s")
    return {"query": state["query"], "ocr_chunks": ocr_chunks}

def answer_synthesis_node(state: Dict) -> Dict[str, Any]:
    import openai
    import time
    config = WineRAGConfig()
    query = state["query"]
    context_parts = []
    for chunk in state.get("retrieved_chunks", []):
        context_parts.append(f"- {chunk[:200]}...")
    for chunk in state.get("vision_chunks", []):
        context_parts.append(f"- {chunk[:200]}...")
    for chunk in state.get("ocr_chunks", []):
        context_parts.append(f"- {chunk[:200]}...")
    context_str = "\n".join(context_parts)
    prompt = f"""
    You are a wine expert. Answer this wine question using ONLY the provided context below. If the answer is not in the context, say 'I don't know based on the provided information.' Do NOT use any outside knowledge.
    \nQuestion: "{query}"
    \nWine Knowledge Context:\n{context_str}
    Provide a helpful, detailed response that:
    1. Directly answers the question using ONLY the context
    2. Includes specific wine recommendations when appropriate
    3. Mentions food pairings if relevant
    4. Suggests alternatives if asked
    5. Uses a friendly, knowledgeable tone
    6. If the answer is not in the context, say 'I don't know based on the provided information.'
    Keep the response comprehensive but concise (max 300 words).
    """
    try:
        client = openai.OpenAI(api_key=config.openai_api_key)
        t0 = time.perf_counter()
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
            max_tokens=800
        )
        t1 = time.perf_counter()
        logging.info(f"[PROFILE] OpenAI completion time: {t1 - t0:.3f}s")
        answer = response.choices[0].message.content
    except Exception as e:
        answer = "I apologize, but I'm having trouble generating a response right now. Please try again."
    no_answer = not (state.get("retrieved_chunks") or state.get("vision_chunks") or state.get("ocr_chunks"))
    next_node = "fallback" if no_answer else None
    return {
        "no_answer": no_answer,
        "final_response": answer,
        "__next__": next_node
    }

def fallback_node(state: Dict) -> Dict[str, Any]:
    return {"final_response": "Sorry, no answer could be generated from the available data."}

# --- Agent Graph Construction ---
def build_agent_graph():
    workflow = StateGraph(dict)
    workflow.add_node("analyze_query", analyze_query_node)
    workflow.add_node("vector_search", vector_search_node)
    workflow.add_node("vision_analysis", vision_analysis_node)
    workflow.add_node("ocr_analysis", ocr_analysis_node)
    workflow.add_node("answer_synthesis", answer_synthesis_node)
    workflow.add_node("fallback", fallback_node)
    workflow.set_entry_point("analyze_query")
    workflow.add_edge("analyze_query", "vector_search")
    workflow.add_edge("vision_analysis", "ocr_analysis")
    workflow.add_edge("ocr_analysis", "answer_synthesis")
    workflow.add_edge("fallback", END)
    # Conditional/dynamic transitions using add_conditional_edges
    workflow.add_conditional_edges("vector_search", lambda state: state.get("__next__"), {
        "answer_synthesis": "answer_synthesis",
        "vision_analysis": "vision_analysis"
    })
    workflow.add_conditional_edges("vision_analysis", lambda state: state.get("__next__"), {
        "answer_synthesis": "answer_synthesis",
        "ocr_analysis": "ocr_analysis"
    })
    workflow.add_conditional_edges("answer_synthesis", lambda state: state.get("__next__"), {
        "fallback": "fallback",
        None: END
    })
    return workflow.compile()
